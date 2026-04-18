"""
Génère le PowerPoint de présentation du projet JumiaPrix CI.

Usage :
    pip install python-pptx
    python generate_pptx.py

Sortie :
    presentation_jumiaprix.pptx (dans le dossier courant)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ────────────────────────────────────────────────────────────
# Palette
# ────────────────────────────────────────────────────────────
BLEU_NUIT   = RGBColor(0x1E, 0x3A, 0x8A)
BLEU_CLAIR  = RGBColor(0xDB, 0xEA, 0xFE)
ORANGE      = RGBColor(0xF9, 0x73, 0x16)
ORANGE_PALE = RGBColor(0xFF, 0xED, 0xD5)
BLANC       = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_FONCE  = RGBColor(0x33, 0x33, 0x33)
GRIS_CLAIR  = RGBColor(0xF3, 0xF4, 0xF6)
GRIS_MOYEN  = RGBColor(0x9C, 0xA3, 0xAF)
VERT        = RGBColor(0x16, 0xA3, 0x4A)
VERT_PALE   = RGBColor(0xDC, 0xFC, 0xE7)
ROUGE       = RGBColor(0xDC, 0x26, 0x26)


# ────────────────────────────────────────────────────────────
# Setup présentation 16:9 (33.87 cm x 19.05 cm ≈ Inches(13.333)x7.5)
# ────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]


# ────────────────────────────────────────────────────────────
# Helpers
# ────────────────────────────────────────────────────────────
def add_blank_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=BLEU_NUIT, line=None, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    if not shadow:
        shape.shadow.inherit = False
    return shape


def add_rounded_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text,
             size=18, bold=False, color=GRIS_FONCE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(60000)
    tf.margin_right = Emu(60000)
    tf.margin_top = Emu(40000)
    tf.margin_bottom = Emu(40000)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=18, color=GRIS_FONCE, bold_first=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(80000)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold_first and i == 0
    return tb


def add_title_bar(slide, title, subtitle=None):
    """Barre de titre standard en haut de slide."""
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), fill=BLEU_NUIT)
    add_rect(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.08), fill=ORANGE)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.9),
             title, size=32, bold=True, color=BLANC,
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.4),
                 subtitle, size=14, color=ORANGE_PALE,
                 anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, page_num, total=16):
    add_rect(slide, Inches(0), Inches(7.15), Inches(13.333), Inches(0.35), fill=GRIS_CLAIR)
    add_text(slide, Inches(0.4), Inches(7.18), Inches(6), Inches(0.3),
             "JumiaPrix CI — ENSEA AS Data Science",
             size=10, color=GRIS_MOYEN)
    add_text(slide, Inches(11), Inches(7.18), Inches(1.9), Inches(0.3),
             f"{page_num} / {total}",
             size=10, color=GRIS_MOYEN, align=PP_ALIGN.RIGHT)


def add_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes


# ────────────────────────────────────────────────────────────
# SLIDE 1 — Page de garde
# ────────────────────────────────────────────────────────────
s = add_blank_slide()
set_bg(s, BLEU_NUIT)

# Bande orange décorative
add_rect(s, Inches(0), Inches(3.2), Inches(13.333), Inches(0.1), fill=ORANGE)

# Titre principal
add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.2),
         "JumiaPrix CI",
         size=72, bold=True, color=BLANC,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Sous-titre
add_text(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.6),
         "Comparateur de prix e-commerce en Côte d'Ivoire",
         size=24, color=ORANGE_PALE,
         align=PP_ALIGN.CENTER)

# Description
add_text(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.6),
         "Pipeline de web scraping multi-source avec IA, alertes et monitoring",
         size=16, color=BLEU_CLAIR,
         align=PP_ALIGN.CENTER)

# Informations école / auteur (bas de page)
add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.4),
         "Nathan KINDO",
         size=18, bold=True, color=BLANC,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.35),
         "ENSEA — AS Data Science — Module Web Scraping",
         size=14, color=BLEU_CLAIR,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.35),
         "2026",
         size=12, color=GRIS_MOYEN,
         align=PP_ALIGN.CENTER)

add_notes(s, "Bonjour, je vais vous présenter mon projet de module Web Scraping : "
             "JumiaPrix CI, un comparateur de prix e-commerce pour la Côte d'Ivoire. "
             "Le projet agrège 3 sources marchandes en ligne et expose les données "
             "via une API, un site web et un chatbot. Durée de présentation : environ 12 minutes.")


# ────────────────────────────────────────────────────────────
# SLIDE 2 — Sommaire
# ────────────────────────────────────────────────────────────
s = add_blank_slide()
set_bg(s, BLANC)
add_title_bar(s, "Sommaire")

sommaire_items = [
    ("1", "Problématique et objectifs"),
    ("2", "Démonstration fonctionnelle"),
    ("3", "Architecture technique (12 services Docker)"),
    ("4", "Pipeline de données et orchestration"),
    ("5", "Base de données (schéma 8 tables)"),
    ("6", "Défis techniques rencontrés"),
    ("7", "Chatbot JumiBot (IA hybride)"),
    ("8", "Monitoring, tests et éthique"),
    ("9", "Chiffres-clés et conclusion"),
]

y = 1.7
for num, text in sommaire_items:
    add_rounded_rect(s, Inches(1.5), Inches(y), Inches(0.7), Inches(0.55), fill=ORANGE)
    add_text(s, Inches(1.5), Inches(y), Inches(0.7), Inches(0.55),
             num, size=20, bold=True, color=BLANC,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.4), Inches(y), Inches(10), Inches(0.55),
             text, size=20, color=GRIS_FONCE, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.58

add_footer(s, 2)
add_notes(s, "Voici le plan de ma présentation. Je commencerai par le problème que j'ai "
             "voulu résoudre, puis je vous montrerai concrètement ce que fait le projet. "
             "Ensuite je détaillerai l'architecture technique, la base de données, les "
             "défis rencontrés, et je terminerai par les chiffres-clés.")


# ────────────────────────────────────────────────────────────
# SLIDE 3 — Problématique
# ────────────────────────────────────────────────────────────
s = add_blank_slide()
set_bg(s, BLANC)
add_title_bar(s, "Problématique", "Le constat de départ")

# Bloc contexte
add_rounded_rect(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.2), fill=GRIS_CLAIR)
add_text(s, Inches(0.8), Inches(1.85), Inches(11.7), Inches(1), 
         "En Côte d'Ivoire, les consommateurs doivent visiter manuellement "
         "plusieurs sites e-commerce pour comparer les prix, repérer les promotions "
         "et suivre les baisses.",
         size=18, color=GRIS_FONCE, anchor=MSO_ANCHOR.MIDDLE)

# 3 colonnes "Problèmes"
problems = [
    ("Temps perdu", "Visite manuelle de 3+ sites"),
    ("Pas d'historique", "Impossible de suivre les tendances de prix"),
    ("Pas d'alerte", "On rate les baisses importantes"),
]
col_w = 3.8
col_start_x = 0.6
col_y = 3.3
for i, (titre, desc) in enumerate(problems):
    x = col_start_x + i * (col_w + 0.3)
    add_rounded_rect(s, Inches(x), Inches(col_y), Inches(col_w), Inches(2.2),
                     fill=BLANC, line=ROUGE)
    add_rounded_rect(s, Inches(x), Inches(col_y), Inches(col_w), Inches(0.7), fill=ROUGE)
    add_text(s, Inches(x), Inches(col_y), Inches(col_w), Inches(0.7),
             titre, size=18, bold=True, color=BLANC,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(x), Inches(col_y + 0.85), Inches(col_w), Inches(1.2),
             desc, size=14, color=GRIS_FONCE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Flèche vers solution
add_text(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
         "→  Besoin d'un service qui AGRÈGE, STOCKE et NOTIFIE automatiquement",
         size=20, bold=True, color=ORANGE,
         align=PP_ALIGN.CENTER)

add_footer(s, 3)
add_notes(s, "Le constat est simple : un acheteur ivoirien qui veut comparer les prix "
             "doit aller sur Jumia, DjokStore, CoinAfrique, et parfois d'autres. "
             "Il perd du temps, n'a aucun historique de prix pour savoir si c'est "
             "vraiment une bonne affaire, et ne peut pas être alerté automatiquement. "
             "L'idée du projet est donc de créer un service qui fait tout ça à sa place.")


# ────────────────────────────────────────────────────────────
# SLIDE 4 — Objectifs
# ────────────────────────────────────────────────────────────
s = add_blank_slide()
set_bg(s, BLANC)
add_title_bar(s, "Objectifs", "Ce que JumiaPrix CI réalise")

objectifs = [
    ("Scraper", "3 sources e-commerce ivoiriennes quotidiennement\nJumia, DjokStore, CoinAfrique"),
    ("Nettoyer", "Normaliser, valider et dédupliquer les données\nTaux de rétention 85-100 %"),
    ("Stocker", "Historiser les prix dans PostgreSQL\n13 600 snapshots à ce jour"),
    ("Exposer", "API REST + interface web + chatbot LLM\nRecherche, filtres, stats, auth"),
    ("Alerter", "Notifications email sur baisses de prix\nPersonnalisables par utilisateur"),
    ("Monitorer", "Prometheus + Grafana + logs structurés\nTemps réel, dashboard visuel"),
]

# Grille 3x2
for i, (titre, desc) in enumerate(objectifs):
    row, col = divmod(i, 3)
    x = 0.5 + col * 4.3
    y = 1.6 + row * 2.5
    add_rounded_rect(s, Inches(x), Inches(y), Inches(4.1), Inches(2.3),
                     fill=BLEU_CLAIR)
    add_rounded_rect(s, Inches(x), Inches(y), Inches(4.1), Inches(0.7), fill=BLEU_NUIT)
    add_text(s, Inches(x), Inches(y), Inches(4.1), Inches(0.7),
             titre, size=20, bold=True, color=BLANC,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(x + 0.2), Inches(y + 0.85), Inches(3.7), Inches(1.35),
             desc, size=14, color=GRIS_FONCE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 4)
add_notes(s, "Six objectifs concrets, que l'on peut résumer par 6 verbes : "
             "scraper, nettoyer, stocker, exposer, alerter, monitorer. "
             "C'est un pipeline complet de la donnée brute jusqu'à l'utilisateur final. "
             "Je détaillerai chaque brique dans les slides suivantes.")


# ────────────────────────────────────────────────────────────
# SLIDE 5 — Démonstration fonctionnelle
# ────────────────────────────────────────────────────────────
s = add_blank_slide()
set_bg(s, BLANC)
add_title_bar(s, "Démonstration fonctionnelle", "Ce que voit l'utilisateur")

pages = [
    ("Page d'accueil", "Vitrine, stats globales, derniers produits"),
    ("Boutique", "Catalogue filtrable par source, catégorie, prix"),
    ("Comparer", "Vue côte à côte entre les 3 sources"),
    ("Historique", "Graphique d'évolution du prix d'un produit"),
    ("Mon espace", "Favoris, alertes, préférences (auth JWT)"),
    ("Assistant IA", "Chatbot conversationnel multilingue"),
    ("Admin", "Gestion users, scrape logs, stats admin"),
    ("API Docs", "Documentation Swagger interactive"),
]

# Grille 4x2 de cartes
for i, (titre, desc) in enumerate(pages):
    row, col = divmod(i, 4)
    x = 0.4 + col * 3.22
    y = 1.6 + row * 2.6
    # Fond navigateur-like
    add_rounded_rect(s, Inches(x), Inches(y), Inches(3.1), Inches(2.3),
                     fill=BLANC, line=GRIS_MOYEN)
    # Barre d'adresse
    add_rect(s, Inches(x), Inches(y), Inches(3.1), Inches(0.35), fill=GRIS_CLAIR)
    # 3 pastilles fenêtre
    for j, col_dot in enumerate([RGBColor(0xEF, 0x44, 0x44),
                                 RGBColor(0xFA, 0xCC, 0x15),
                                 RGBColor(0x22, 0xC5, 0x5E)]):
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x + 0.1 + j * 0.18), Inches(y + 0.08),
                                 Inches(0.14), Inches(0.14))
        dot.fill.solid()
        dot.fill.fore_color.rgb = col_dot
        dot.line.fill.background()
    # Titre page
    add_rect(s, Inches(x), Inches(y + 0.35), Inches(3.1), Inches(0.7), fill=BLEU_NUIT)
    add_text(s, Inches(x), Inches(y + 0.35), Inches(3.1), Inches(0.7),
             titre, size=15, bold=True, color=BLANC,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Description
    add_text(s, Inches(x + 0.15), Inches(y + 1.15), Inches(2.8), Inches(1.1),
             desc, size=11, color=GRIS_FONCE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 5)
add_notes(s, "Côté utilisateur, le projet propose 8 pages. "
             "La boutique permet de filtrer le catalogue, "
             "la page Comparer met les 3 sources côte à côte, "
             "Historique affiche un graphique d'évolution du prix, "
             "Mon Espace gère les favoris et les alertes. "
             "Et un chatbot IA répond aux questions en langage naturel. "
             "Tout ça est protégé par une authentification JWT.")


# ────────────────────────────────────────────────────────────
# SLIDE 6 — Architecture globale
# ────────────────────────────────────────────────────────────
s = add_blank_slide()
set_bg(s, BLANC)
add_title_bar(s, "Architecture technique", "12 services orchestrés par Docker Compose")

# 4 couches horizontales
couches = [
    ("SOURCES", [("Jumia CI", ORANGE), ("DjokStore", ORANGE), ("CoinAfrique", ORANGE),
                 ("FlareSolverr", GRIS_MOYEN)], 1.6),
    ("COLLECTE & TRAITEMENT", [("Scrapy spiders", VERT), ("Celery Worker", VERT),
                               ("Pandas cleaner", VERT), ("Celery Beat", VERT)], 2.8),
    ("STOCKAGE & API", [("PostgreSQL", BLEU_NUIT), ("Redis", BLEU_NUIT),
                        ("Flask API", BLEU_NUIT), ("Airflow", BLEU_NUIT)], 4.0),
    ("OBSERVABILITÉ & UX", [("Prometheus", RGBColor(0x7C, 0x3A, 0xED)),
                            ("Grafana", RGBColor(0x7C, 0x3A, 0xED)),
                            ("Chatbot Groq", RGBColor(0x7C, 0x3A, 0xED)),
                            ("Web UI", RGBColor(0x7C, 0x3A, 0xED))], 5.2),
]

for label, boxes, y in couches:
    # Label couche à gauche
    add_text(s, Inches(0.3), Inches(y + 0.15), Inches(1.6), Inches(0.8),
             label, size=12, bold=True, color=GRIS_MOYEN,
             anchor=MSO_ANCHOR.MIDDLE)
    # Boîtes
    for i, (name, color) in enumerate(boxes):
        x = 2.0 + i * 2.75
        add_rounded_rect(s, Inches(x), Inches(y), Inches(2.6), Inches(0.9),
                         fill=color)
        add_text(s, Inches(x), Inches(y), Inches(2.6), Inches(0.9),
                 name, size=14, bold=True, color=BLANC,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Flèches verticales entre couches (stylisées)
for y_arrow in [2.5, 3.7, 4.9]:
    for x_arrow in [3.0, 5.75, 8.5, 11.25]:
        conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x_arrow), Inches(y_arrow),
                                      Inches(x_arrow), Inches(y_arrow + 0.3))
        conn.line.color.rgb = GRIS_MOYEN
        conn.line.width = Pt(1.5)

# Légende bas
add_rounded_rect(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.75),
                 fill=GRIS_CLAIR)
add_text(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.65),
         "Tout est containerisé. Un seul « docker compose up » démarre l'ensemble. "
         "Les services communiquent via un réseau Docker interne (noms de services, pas localhost).",
         size=13, color=GRIS_FONCE, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 6)
add_notes(s, "L'architecture est organisée en 4 couches : "
             "les sources que l'on scrape, la couche de collecte et traitement avec Scrapy et Celery, "
             "le stockage et l'API avec PostgreSQL, Redis et Flask, "
             "et enfin la couche d'observabilité avec Prometheus, Grafana et le chatbot. "
             "Au total 12 services Docker, orchestrés par Docker Compose. Un seul 'docker compose up' "
             "démarre tout l'environnement.")


# ────────────────────────────────────────────────────────────
# SLIDE 7 — Stack technologique
# ────────────────────────────────────────────────────────────
s = add_blank_slide()
set_bg(s, BLANC)
add_title_bar(s, "Stack technologique", "Technologies utilisées par domaine")

stack = [
    ("Langage", "Python 3.12"),
    ("Web scraping", "Scrapy 2.12, Parsel, FlareSolverr"),
    ("Nettoyage", "Pandas 2.2.3"),
    ("Base de données", "PostgreSQL 16, SQLAlchemy 2, Alembic"),
    ("API web", "Flask 3, Flasgger (Swagger), Jinja2"),
    ("Authentification", "Flask-JWT-Extended, bcrypt"),
    ("Tâches async", "Celery 5.4, Redis 7"),
    ("Orchestration", "Apache Airflow 2.9"),
    ("Intelligence artificielle", "Groq API + Llama 3.3 70B"),
    ("Monitoring", "Prometheus, Grafana"),
    ("Tests", "pytest (42 tests)"),
    ("Infrastructure", "Docker, Docker Compose"),
]

# Tableau 2 colonnes x 6 lignes
row_h = 0.82
header_y = 1.5
add_rect(s, Inches(0.5), Inches(header_y), Inches(12.3), Inches(0.6), fill=BLEU_NUIT)
add_text(s, Inches(0.8), Inches(header_y), Inches(4.2), Inches(0.6),
         "DOMAINE", size=14, bold=True, color=BLANC, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(5.2), Inches(header_y), Inches(7.5), Inches(0.6),
         "TECHNOLOGIES", size=14, bold=True, color=BLANC, anchor=MSO_ANCHOR.MIDDLE)

for i, (dom, tech) in enumerate(stack):
    y = header_y + 0.6 + i * 0.38
    bg = GRIS_CLAIR if i % 2 == 0 else BLANC
    add_rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(0.38), fill=bg)
    add_text(s, Inches(0.8), Inches(y), Inches(4.2), Inches(0.38),
             dom, size=12, bold=True, color=BLEU_NUIT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.2), Inches(y), Inches(7.5), Inches(0.38),
             tech, size=12, color=GRIS_FONCE, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 7)
add_notes(s, "Le projet utilise exclusivement des outils open-source éprouvés en production. "
             "Python 3.12 comme langage principal. Scrapy pour le scraping car c'est un framework "
             "robuste avec gestion native des délais et des retries. "
             "PostgreSQL pour sa fiabilité transactionnelle. "
             "Flask pour sa simplicité et sa documentation. "
             "Celery pour l'asynchrone. Airflow pour les pipelines complexes. "
             "Et pour l'IA, Llama 3.3 70B via l'API Groq, qui a un très bon rapport qualité / latence.")


# ────────────────────────────────────────────────────────────
# SLIDE 8 — Pipeline de données (Airflow DAG)
# ────────────────────────────────────────────────────────────
s = add_blank_slide()
set_bg(s, BLANC)
add_title_bar(s, "Pipeline de données", "DAG Airflow — exécution quotidienne")

# Définition des nœuds du DAG
nodes = [
    ("scrape_jumia",       1.0, 2.5, ORANGE,    "Scrapy + FlareSolverr"),
    ("scrape_djokstore",   1.0, 3.8, ORANGE,    "Scrapy"),
    ("scrape_coinafrique", 1.0, 5.1, ORANGE,    "Scrapy"),
    ("clean_and_insert",   4.5, 3.8, VERT,      "Pandas + SQLAlchemy"),
    ("check_price_drops",  7.5, 3.8, BLEU_NUIT, "Détection baisses"),
    ("check_price_alerts", 10.0, 2.8, BLEU_NUIT,"Email SMTP"),
    ("match_cross_source", 10.0, 4.8, BLEU_NUIT,"Fuzzy matching"),
]

node_w, node_h = 2.5, 0.9
for name, x, y, color, subtitle in nodes:
    add_rounded_rect(s, Inches(x), Inches(y), Inches(node_w), Inches(node_h),
                     fill=color)
    add_text(s, Inches(x), Inches(y + 0.05), Inches(node_w), Inches(0.5),
             name, size=13, bold=True, color=BLANC,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(x), Inches(y + 0.5), Inches(node_w), Inches(0.4),
             subtitle, size=10, color=BLANC,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Flèches entre nœuds (connecteurs droits)
arrows = [
    (1.0 + node_w, 2.5 + node_h / 2, 4.5, 3.8 + node_h / 2),  # jumia → clean
    (1.0 + node_w, 3.8 + node_h / 2, 4.5, 3.8 + node_h / 2),  # djokstore → clean
    (1.0 + node_w, 5.1 + node_h / 2, 4.5, 3.8 + node_h / 2),  # coin → clean
    (4.5 + node_w, 3.8 + node_h / 2, 7.5, 3.8 + node_h / 2),  # clean → drops
    (7.5 + node_w, 3.8 + node_h / 2, 10.0, 2.8 + node_h / 2), # drops → alerts
    (7.5 + node_w, 3.8 + node_h / 2, 10.0, 4.8 + node_h / 2), # drops → match
]
for x1, y1, x2, y2 in arrows:
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                  Inches(x1), Inches(y1),
                                  Inches(x2), Inches(y2))
    conn.line.color.rgb = GRIS_MOYEN
    conn.line.width = Pt(2)

# Légendes scheduling
add_rounded_rect(s, Inches(0.5), Inches(6.4), Inches(6), Inches(0.65), fill=BLEU_CLAIR)
add_text(s, Inches(0.6), Inches(6.4), Inches(5.8), Inches(0.65),
         "Airflow : scraping quotidien + weekly_digest (lundi 8 h)",
         size=13, bold=True, color=BLEU_NUIT, anchor=MSO_ANCHOR.MIDDLE)

add_rounded_rect(s, Inches(6.8), Inches(6.4), Inches(6), Inches(0.65), fill=ORANGE_PALE)
add_text(s, Inches(6.9), Inches(6.4), Inches(5.8), Inches(0.65),
         "Celery Beat : check chutes majeures toutes les 5 min",
         size=13, bold=True, color=ORANGE, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 8)
add_notes(s, "Voici le pipeline quotidien orchestré par Airflow. "
             "Les 3 spiders peuvent tourner en parallèle, ils se rejoignent au nettoyage "
             "puis à l'insertion en base. Ensuite on détecte les baisses de prix, "
             "ce qui déclenche deux tâches en parallèle : l'envoi des alertes par email "
             "et le matching cross-source, qui rapproche un même produit vendu par plusieurs "
             "sites grâce à du fuzzy matching. "
             "Airflow gère tout ce qui a des dépendances. Celery Beat, lui, ne garde qu'une seule "
             "tâche : une vérification très fréquente, toutes les 5 minutes, des chutes majeures. "
             "C'est une séparation claire des responsabilités.")


# ────────────────────────────────────────────────────────────
# SLIDE 9 — Schéma de la base de données
# ────────────────────────────────────────────────────────────
s = add_blank_slide()
set_bg(s, BLANC)
add_title_bar(s, "Base de données", "8 tables PostgreSQL avec contraintes et vues")


def draw_table(slide, x, y, w, h, name, cols, fill_header=BLEU_NUIT, fill_body=BLEU_CLAIR):
    # En-tête
    add_rect(slide, Inches(x), Inches(y), Inches(w), Inches(0.4), fill=fill_header)
    add_text(slide, Inches(x), Inches(y), Inches(w), Inches(0.4),
             name, size=12, bold=True, color=BLANC,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Corps
    add_rect(slide, Inches(x), Inches(y + 0.4), Inches(w), Inches(h - 0.4),
             fill=fill_body, line=fill_header)
    tb = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 0.42),
                                  Inches(w - 0.1), Inches(h - 0.45))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, col in enumerate(cols):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = col
        run.font.name = "Consolas"
        run.font.size = Pt(9)
        run.font.color.rgb = GRIS_FONCE


# Couleurs par catégorie
USER_H, USER_B = VERT, VERT_PALE
DATA_H, DATA_B = BLEU_NUIT, BLEU_CLAIR
LOG_H, LOG_B   = GRIS_MOYEN, GRIS_CLAIR

# Position X: 0.5 - 13 (disponible = 12.5"), Y: 1.4 - 7 (disponible = 5.6")

# Table products (CENTRE)
draw_table(s, 5.3, 3.1, 2.7, 2.0, "products",
           ["PK id", "source", "source_product_id", "name", "category",
            "image_url", "product_url", "currency", "UK(source, src_id)"],
           DATA_H, DATA_B)

# price_history (DROITE) - plus grande
draw_table(s, 9.0, 3.1, 2.6, 2.0, "price_history",
           ["PK id", "FK product_id", "price", "old_price",
            "discount_pct", "reviews_count", "scraped_at",
            "CHECK price > 0",
            "13 600 lignes"],
           ORANGE, ORANGE_PALE)

# product_matches (HAUT DROITE)
draw_table(s, 9.5, 1.4, 2.3, 1.3, "product_matches",
           ["PK id", "FK product_a_id", "FK product_b_id",
            "similarity_score", "matched_at"],
           DATA_H, DATA_B)

# users (GAUCHE)
draw_table(s, 0.5, 1.4, 2.5, 1.5, "users",
           ["PK id", "email UNIQUE", "password_hash", "username",
            "is_admin", "created_at"],
           USER_H, USER_B)

# user_preferences (SOUS users)
draw_table(s, 0.5, 3.1, 2.5, 1.2, "user_preferences",
           ["PK id", "FK user_id", "categories_preferees[]",
            "budget_max", "sources_preferees[]"],
           USER_H, USER_B)

# user_favorites (MILIEU GAUCHE)
draw_table(s, 0.5, 4.5, 2.5, 1.0, "user_favorites",
           ["PK id", "FK user_id", "FK product_id", "created_at"],
           USER_H, USER_B)

# price_alerts (BAS GAUCHE)
draw_table(s, 0.5, 5.7, 2.5, 1.2, "price_alerts",
           ["PK id", "FK user_id", "FK product_id", "target_price",
            "active", "last_notified_at"],
           USER_H, USER_B)

# scrape_logs (BAS DROITE, isolée)
draw_table(s, 9.0, 5.7, 2.6, 1.2, "scrape_logs",
           ["PK id", "source", "status", "items_raw",
            "items_clean", "duration_sec", "started_at"],
           LOG_H, LOG_B)


def add_fk_line(slide, x1, y1, x2, y2, color=ORANGE):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.5)


# Relations
# users → user_preferences
add_fk_line(s, 1.75, 2.9, 1.75, 3.1, VERT)
# users → user_favorites
add_fk_line(s, 1.75, 2.9, 1.75, 4.5, VERT)
# users → price_alerts
add_fk_line(s, 1.75, 2.9, 1.75, 5.7, VERT)
# user_favorites → products
add_fk_line(s, 3.0, 5.0, 5.3, 4.1, ORANGE)
# price_alerts → products
add_fk_line(s, 3.0, 6.3, 5.3, 4.3, ORANGE)
# products → price_history
add_fk_line(s, 8.0, 4.1, 9.0, 4.1, ORANGE)
# products → product_matches (self join)
add_fk_line(s, 6.6, 3.1, 10.6, 2.7, ORANGE)

# Légende
add_rounded_rect(s, Inches(3.4), Inches(1.4), Inches(5.7), Inches(1.5), fill=GRIS_CLAIR)
add_text(s, Inches(3.55), Inches(1.5), Inches(5.5), Inches(0.4),
         "Caractéristiques du schéma", size=13, bold=True, color=BLEU_NUIT)
legend_items = [
    "• 8 tables, contraintes CHECK, UNIQUE, FK CASCADE",
    "• Indexes composites sur (product_id, scraped_at)",
    "• Triggers plpgsql pour updated_at",
    "• 2 vues : v_latest_prices (JOIN LATERAL), v_price_evolution",
]
add_bullets(s, Inches(3.55), Inches(1.9), Inches(5.5), Inches(1.1),
            legend_items, size=10)

# Codes couleurs
legend_y = 6.95
add_rect(s, Inches(0.5), Inches(legend_y), Inches(0.3), Inches(0.15), fill=VERT)
add_text(s, Inches(0.85), Inches(legend_y - 0.02), Inches(2.5), Inches(0.2),
         "Tables utilisateur", size=10, color=GRIS_FONCE)
add_rect(s, Inches(3.5), Inches(legend_y), Inches(0.3), Inches(0.15), fill=BLEU_NUIT)
add_text(s, Inches(3.85), Inches(legend_y - 0.02), Inches(2.5), Inches(0.2),
         "Tables data", size=10, color=GRIS_FONCE)
add_rect(s, Inches(6.0), Inches(legend_y), Inches(0.3), Inches(0.15), fill=ORANGE)
add_text(s, Inches(6.35), Inches(legend_y - 0.02), Inches(2.5), Inches(0.2),
         "Historique (volumineuse)", size=10, color=GRIS_FONCE)
add_rect(s, Inches(9.0), Inches(legend_y), Inches(0.3), Inches(0.15), fill=GRIS_MOYEN)
add_text(s, Inches(9.35), Inches(legend_y - 0.02), Inches(2.5), Inches(0.2),
         "Logs", size=10, color=GRIS_FONCE)

add_footer(s, 9)
add_notes(s, "Le schéma de la base contient 8 tables organisées en 3 groupes. "
             "En vert, les tables utilisateur : users, user_preferences, user_favorites, price_alerts. "
             "En bleu, les tables de données : products au centre comme table pivot, product_matches pour "
             "le matching cross-source. "
             "En orange, la table price_history, la plus volumineuse avec plus de 13 600 lignes "
             "d'historique de prix. "
             "En gris, scrape_logs, isolée, qui journalise chaque exécution. "
             "Les relations avec clés étrangères sont en CASCADE pour garantir la cohérence. "
             "J'ai aussi créé 2 vues SQL, dont une avec un JOIN LATERAL pour récupérer efficacement "
             "le dernier prix de chaque produit.")


# ────────────────────────────────────────────────────────────
# SLIDE 10 — Défi : Cloudflare sur Jumia
# ────────────────────────────────────────────────────────────
s = add_blank_slide()
set_bg(s, BLANC)
add_title_bar(s, "Défi technique : Cloudflare Turnstile", "Comment scraper Jumia malgré la protection anti-bot")

# Problème
add_rounded_rect(s, Inches(0.5), Inches(1.6), Inches(6.1), Inches(2.6),
                 fill=BLANC, line=ROUGE)
add_rect(s, Inches(0.5), Inches(1.6), Inches(6.1), Inches(0.5), fill=ROUGE)
add_text(s, Inches(0.5), Inches(1.6), Inches(6.1), Inches(0.5),
         "PROBLÈME", size=14, bold=True, color=BLANC,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
problem_items = [
    "Jumia utilise Cloudflare Turnstile (CAPTCHA JS)",
    "Scrapy pur = erreur 403 Forbidden immédiate",
    "Les headers et User-Agent ne suffisent pas",
    "Un token dynamique est exigé à chaque requête",
]
add_bullets(s, Inches(0.7), Inches(2.25), Inches(5.8), Inches(1.9),
            problem_items, size=13)

# Solution
add_rounded_rect(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(2.6),
                 fill=BLANC, line=VERT)
add_rect(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.5), fill=VERT)
add_text(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.5),
         "SOLUTION", size=14, bold=True, color=BLANC,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
solution_items = [
    "FlareSolverr : proxy Chromium headless",
    "Le spider passe via FlareSolverr qui résout Turnstile",
    "Retry Celery (max_retries=2) en cas d'échec",
    "Timeout configuré pour éviter les blocages",
]
add_bullets(s, Inches(7.0), Inches(2.25), Inches(5.7), Inches(1.9),
            solution_items, size=13)

# Preuve en prod (flux illustré)
add_rounded_rect(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.3),
                 fill=GRIS_CLAIR)
add_text(s, Inches(0.7), Inches(4.7), Inches(12), Inches(0.4),
         "Preuve en production (logs réels)",
         size=14, bold=True, color=BLEU_NUIT)

# Ligne de tentatives
attempts = [
    ("08:11", "❌ Error",   ROUGE),
    ("08:12", "❌ Error",   ROUGE),
    ("08:13", "❌ Error",   ROUGE),
    ("08:14", "✅ Success", VERT),
    ("",      "780 items",  BLEU_NUIT),
]
for i, (time, status, color) in enumerate(attempts):
    x = 0.8 + i * 2.4
    add_rounded_rect(s, Inches(x), Inches(5.3), Inches(2.2), Inches(1.3), fill=BLANC, line=color)
    add_text(s, Inches(x), Inches(5.4), Inches(2.2), Inches(0.4),
             time, size=11, bold=True, color=GRIS_MOYEN,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(x), Inches(5.8), Inches(2.2), Inches(0.6),
             status, size=16, bold=True, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.25),
         "→ Le pipeline retry automatiquement et finit par passer. Résilience prouvée.",
         size=12, color=GRIS_FONCE, align=PP_ALIGN.CENTER)

add_footer(s, 10)
add_notes(s, "Le principal défi technique était Cloudflare Turnstile sur Jumia. "
             "Scrapy seul reçoit un 403 immédiat. J'ai intégré FlareSolverr, un proxy "
             "basé sur Chromium headless qui résout le challenge JavaScript. "
             "Le spider envoie sa requête à FlareSolverr qui retourne le HTML une fois le "
             "token obtenu. En plus, j'ai activé 2 retries dans Celery. "
             "La preuve que ça marche : sur un run récent en production, j'ai eu 3 erreurs "
             "consécutives à 8h11, 8h12, 8h13, puis un succès à 8h14 avec 780 items. "
             "Le système s'auto-répare sans intervention.")


# ────────────────────────────────────────────────────────────
# SLIDE 11 — Chatbot JumiBot
# ────────────────────────────────────────────────────────────
s = add_blank_slide()
set_bg(s, BLANC)
add_title_bar(s, "Chatbot JumiBot", "Architecture hybride IA + logique locale")

# Schéma : user → API → 2 paths
# User
add_rounded_rect(s, Inches(0.5), Inches(3.2), Inches(2.0), Inches(1.0), fill=GRIS_MOYEN)
add_text(s, Inches(0.5), Inches(3.2), Inches(2.0), Inches(1.0),
         "Utilisateur", size=14, bold=True, color=BLANC,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Flèche
conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                              Inches(2.5), Inches(3.7),
                              Inches(3.5), Inches(3.7))
conn.line.color.rgb = GRIS_MOYEN
conn.line.width = Pt(2.5)

# API /chat
add_rounded_rect(s, Inches(3.5), Inches(3.2), Inches(2.3), Inches(1.0), fill=BLEU_NUIT)
add_text(s, Inches(3.5), Inches(3.2), Inches(2.3), Inches(1.0),
         "API /chat\n(Flask)", size=13, bold=True, color=BLANC,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Flèches vers 2 chemins
conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                              Inches(5.8), Inches(3.5),
                              Inches(7.2), Inches(1.9))
conn.line.color.rgb = VERT
conn.line.width = Pt(2)
conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                              Inches(5.8), Inches(3.9),
                              Inches(7.2), Inches(5.5))
conn.line.color.rgb = ORANGE
conn.line.width = Pt(2)

# Chemin 1 : NLU local
add_rounded_rect(s, Inches(7.2), Inches(1.4), Inches(5.6), Inches(2.0),
                 fill=VERT_PALE, line=VERT)
add_rect(s, Inches(7.2), Inches(1.4), Inches(5.6), Inches(0.5), fill=VERT)
add_text(s, Inches(7.2), Inches(1.4), Inches(5.6), Inches(0.5),
         "NLU locale (côté serveur)", size=13, bold=True, color=BLANC,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(7.4), Inches(2.0), Inches(5.3), Inches(1.35),
            [
                "Extraction filtres (source, prix, catégorie)",
                "Recherche produits pertinents en DB",
                "Fallback conversationnel (salut, merci...)",
                "Fonctionne sans internet",
            ], size=11)

# Chemin 2 : LLM Groq
add_rounded_rect(s, Inches(7.2), Inches(5.0), Inches(5.6), Inches(2.0),
                 fill=ORANGE_PALE, line=ORANGE)
add_rect(s, Inches(7.2), Inches(5.0), Inches(5.6), Inches(0.5), fill=ORANGE)
add_text(s, Inches(7.2), Inches(5.0), Inches(5.6), Inches(0.5),
         "LLM Llama 3.3 70B (Groq API)", size=13, bold=True, color=BLANC,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(7.4), Inches(5.6), Inches(5.3), Inches(1.35),
            [
                "Reçoit stats DB réelles dans le prompt",
                "Reçoit la liste de produits trouvés",
                "Formule une réponse naturelle et engageante",
                "Zéro hallucination (données factuelles)",
            ], size=11)

# PostgreSQL en bas
add_rounded_rect(s, Inches(0.5), Inches(5.4), Inches(5.8), Inches(1.5),
                 fill=BLEU_CLAIR, line=BLEU_NUIT)
add_rect(s, Inches(0.5), Inches(5.4), Inches(5.8), Inches(0.4), fill=BLEU_NUIT)
add_text(s, Inches(0.5), Inches(5.4), Inches(5.8), Inches(0.4),
         "PostgreSQL — source de vérité", size=12, bold=True, color=BLANC,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.7), Inches(5.9), Inches(5.4), Inches(0.9),
         "Stats (1930 produits, 24 catégories, 1088 promos)\nProduits filtrés selon la requête",
         size=11, color=GRIS_FONCE, anchor=MSO_ANCHOR.MIDDLE)

# Flèche DB → les 2 chemins
conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                              Inches(4.5), Inches(5.4),
                              Inches(7.2), Inches(2.4))
conn.line.color.rgb = BLEU_NUIT
conn.line.width = Pt(1.5)
conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                              Inches(4.5), Inches(5.8),
                              Inches(7.2), Inches(6.0))
conn.line.color.rgb = BLEU_NUIT
conn.line.width = Pt(1.5)

add_footer(s, 11)
add_notes(s, "Le chatbot est hybride, c'est un choix volontaire. "
             "Côté serveur, une couche NLU locale extrait les filtres de la requête "
             "(source, prix, catégorie) avec des regex et des dictionnaires. "
             "Elle gère aussi les salutations et les remerciements localement, "
             "sans appeler l'IA. "
             "Côté IA, j'utilise Llama 3.3 70B via Groq. Mais surtout : "
             "je lui INJECTE dans le prompt les vraies statistiques de la base "
             "et la liste réelle des produits trouvés. "
             "Résultat : le modèle compose des réponses naturelles, mais il ne peut "
             "PAS halluciner car il ne parle que des vrais produits. "
             "Si Groq est indisponible, le fallback local prend le relais.")


# ────────────────────────────────────────────────────────────
# SLIDE 12 — Monitoring
# ────────────────────────────────────────────────────────────
s = add_blank_slide()
set_bg(s, BLANC)
add_title_bar(s, "Monitoring", "Prometheus + Grafana — observabilité en temps réel")

# Prometheus
add_rounded_rect(s, Inches(0.5), Inches(1.6), Inches(6.1), Inches(2.6),
                 fill=BLANC, line=ORANGE)
add_rect(s, Inches(0.5), Inches(1.6), Inches(6.1), Inches(0.5), fill=ORANGE)
add_text(s, Inches(0.5), Inches(1.6), Inches(6.1), Inches(0.5),
         "Prometheus — Collecte", size=14, bold=True, color=BLANC,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
prom_items = [
    "Scrape /metrics toutes les 15 secondes",
    "6 métriques custom exposées par l'API",
    "Stockage en base TSDB intégré",
    "up == 1 sur api:5000 et prometheus:9090",
]
add_bullets(s, Inches(0.7), Inches(2.25), Inches(5.8), Inches(1.9),
            prom_items, size=13)

# Grafana
add_rounded_rect(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(2.6),
                 fill=BLANC, line=BLEU_NUIT)
add_rect(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.5), fill=BLEU_NUIT)
add_text(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.5),
         "Grafana — Visualisation", size=14, bold=True, color=BLANC,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
graf_items = [
    "Dashboard provisionné automatiquement",
    "Graphes : requêtes/sec, latence, erreurs",
    "Auto-refresh 30 s, data source Prometheus",
    "Accessible sur http://localhost:3000",
]
add_bullets(s, Inches(7.0), Inches(2.25), Inches(5.7), Inches(1.9),
            graf_items, size=13)

# Stats réelles de l'API
add_rounded_rect(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.3),
                 fill=GRIS_CLAIR)
add_text(s, Inches(0.7), Inches(4.7), Inches(12), Inches(0.4),
         "Statistiques réelles de production",
         size=14, bold=True, color=BLEU_NUIT)

stats_grid = [
    ("2 546", "GET 200", VERT),
    ("0", "Erreurs 5xx", VERT),
    ("520", "404 (scans bots)", GRIS_MOYEN),
    ("3", "422 (validation)", GRIS_MOYEN),
    ("1", "POST 201", VERT),
]
for i, (val, label, color) in enumerate(stats_grid):
    x = 0.8 + i * 2.4
    add_rounded_rect(s, Inches(x), Inches(5.3), Inches(2.2), Inches(1.3), fill=BLANC)
    add_text(s, Inches(x), Inches(5.4), Inches(2.2), Inches(0.7),
             val, size=30, bold=True, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(x), Inches(6.1), Inches(2.2), Inches(0.4),
             label, size=11, color=GRIS_FONCE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 12)
add_notes(s, "Prometheus scrape l'API toutes les 15 secondes sur l'endpoint /metrics. "
             "Il collecte 6 métriques custom : nombre de requêtes, durée, statut HTTP, endpoint, "
             "méthode, et le nombre de produits en DB. Grafana consomme ces données "
             "pour afficher des dashboards en temps réel. "
             "En production, l'API a servi 2 546 requêtes GET en 200 OK sans la moindre "
             "erreur 5xx. Les 404 correspondent à des scans de bots externes, et les 422 à "
             "des validations utilisateur. Cette stabilité valide les choix d'architecture.")


# ────────────────────────────────────────────────────────────
# SLIDE 13 — Tests et qualité
# ────────────────────────────────────────────────────────────
s = add_blank_slide()
set_bg(s, BLANC)
add_title_bar(s, "Tests et qualité", "Garantir la robustesse du pipeline")

# Grande stat 42
add_rounded_rect(s, Inches(0.5), Inches(1.6), Inches(4.0), Inches(5.3),
                 fill=BLEU_NUIT)
add_text(s, Inches(0.5), Inches(2.3), Inches(4.0), Inches(2.0),
         "42", size=150, bold=True, color=ORANGE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.5), Inches(4.5), Inches(4.0), Inches(0.7),
         "tests unitaires", size=22, bold=True, color=BLANC,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.5), Inches(5.2), Inches(4.0), Inches(0.6),
         "pytest", size=16, color=ORANGE_PALE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(5.9), Inches(4.0), Inches(0.6),
         "> exigence sujet (3)", size=14, color=BLEU_CLAIR,
         align=PP_ALIGN.CENTER)

# Couverture détaillée
categories = [
    ("Validation des prix", "Aberrants, négatifs, None, zéro, hors bornes catégorie"),
    ("Extraction réductions", "Parsing des % depuis texte, cas manquants, cohérence old/new"),
    ("Suppression doublons", "Par URL, par nom normalisé, conservation du plus récent"),
    ("Classification catégories", "Mapping exact, fallback par mots-clés, cas inconnus"),
    ("Pipeline complet", "Intégration scrape → clean → insert sur jeu réel"),
    ("Cas limites", "Données vides, encodages, valeurs NaN, types inattendus"),
]

for i, (titre, desc) in enumerate(categories):
    y = 1.6 + i * 0.88
    add_rounded_rect(s, Inches(4.8), Inches(y), Inches(0.4), Inches(0.75),
                     fill=ORANGE)
    add_text(s, Inches(4.8), Inches(y), Inches(0.4), Inches(0.75),
             f"{i+1}", size=14, bold=True, color=BLANC,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rounded_rect(s, Inches(5.3), Inches(y), Inches(7.5), Inches(0.75),
                     fill=GRIS_CLAIR)
    add_text(s, Inches(5.5), Inches(y + 0.05), Inches(7.2), Inches(0.35),
             titre, size=13, bold=True, color=BLEU_NUIT)
    add_text(s, Inches(5.5), Inches(y + 0.38), Inches(7.2), Inches(0.35),
             desc, size=11, color=GRIS_FONCE)

add_footer(s, 13)
add_notes(s, "La qualité du code est garantie par 42 tests unitaires avec pytest, "
             "ce qui dépasse largement l'exigence minimale de 3 tests du sujet. "
             "Les tests couvrent 6 grandes catégories : "
             "validation des prix avec tous les cas aberrants, extraction des réductions, "
             "suppression de doublons, classification des catégories, pipeline complet "
             "et cas limites. "
             "J'utilise aussi Flasgger qui génère automatiquement une documentation "
             "Swagger/OpenAPI de l'API à partir des docstrings Python.")


# ────────────────────────────────────────────────────────────
# SLIDE 14 — Éthique
# ────────────────────────────────────────────────────────────
s = add_blank_slide()
set_bg(s, BLANC)
add_title_bar(s, "Respect de l'éthique scraping", "Être conforme aux bonnes pratiques")

principles = [
    ("robots.txt", "Respecté. Pour Jumia qui n'expose pas de robots.txt standard, "
                    "enforcement manuel via la configuration Scrapy."),
    ("Délais entre requêtes", "DOWNLOAD_DELAY de 1 à 2 secondes + AUTOTHROTTLE adaptatif "
                               "selon la charge du serveur."),
    ("Volume limité", "CLOSESPIDER_ITEMCOUNT = 500 items par run, conforme à "
                      "la limite du sujet."),
    ("User-Agent identifié", "Pas de spoofing malveillant. L'identité du bot est claire "
                              "dans les logs des sites scrappés."),
    ("Heures de scraping", "Planification via Airflow en dehors des heures de pointe "
                           "pour ne pas dégrader l'expérience utilisateur."),
    ("Données publiques uniquement", "Aucune donnée personnelle collectée. Uniquement "
                                      "catalogue produit accessible sans connexion."),
]

for i, (titre, desc) in enumerate(principles):
    row, col = divmod(i, 2)
    x = 0.5 + col * 6.3
    y = 1.6 + row * 1.75
    add_rounded_rect(s, Inches(x), Inches(y), Inches(6.1), Inches(1.6),
                     fill=VERT_PALE, line=VERT)
    # Checkmark
    add_rounded_rect(s, Inches(x + 0.1), Inches(y + 0.1), Inches(0.6), Inches(0.6),
                     fill=VERT)
    add_text(s, Inches(x + 0.1), Inches(y + 0.1), Inches(0.6), Inches(0.6),
             "✓", size=22, bold=True, color=BLANC,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Titre
    add_text(s, Inches(x + 0.85), Inches(y + 0.1), Inches(5.1), Inches(0.45),
             titre, size=14, bold=True, color=VERT)
    # Description
    add_text(s, Inches(x + 0.85), Inches(y + 0.55), Inches(5.1), Inches(1.0),
             desc, size=11, color=GRIS_FONCE)

add_footer(s, 14)
add_notes(s, "L'éthique est au cœur du projet. "
             "Je respecte robots.txt, avec un enforcement manuel pour Jumia qui n'expose pas "
             "de robots.txt standard. J'ajoute des délais entre requêtes avec AUTOTHROTTLE "
             "qui s'adapte à la charge du serveur. "
             "Je limite le volume à 500 items par run, conformément au sujet. "
             "Le User-Agent identifie clairement le bot. "
             "Je scrape en dehors des heures de pointe. "
             "Et je ne collecte que des données publiques, jamais personnelles.")


# ────────────────────────────────────────────────────────────
# SLIDE 15 — Chiffres-clés
# ────────────────────────────────────────────────────────────
s = add_blank_slide()
set_bg(s, BLEU_NUIT)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.1), fill=BLEU_NUIT)
add_rect(s, Inches(0), Inches(1.1), Inches(13.333), Inches(0.08), fill=ORANGE)
add_text(s, Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.9),
         "Chiffres-clés du projet", size=32, bold=True, color=BLANC,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.4),
         "Résultats concrets en production",
         size=14, color=ORANGE_PALE, anchor=MSO_ANCHOR.MIDDLE)

# Grille 4x2 de chiffres
kpis = [
    ("1 930",  "produits en base",        ORANGE),
    ("13 600", "snapshots de prix",       BLEU_CLAIR),
    ("3",      "sources scrappées",       ORANGE),
    ("24",     "catégories couvertes",    BLEU_CLAIR),
    ("1 088",  "promos actives détectées",ORANGE),
    ("40 %",   "réduction moyenne",       BLEU_CLAIR),
    ("2 546",  "requêtes API servies",    ORANGE),
    ("0",      "erreur 5xx",              BLEU_CLAIR),
    ("42",     "tests pytest",            ORANGE),
    ("12",     "services Docker",         BLEU_CLAIR),
    ("5 000+", "lignes de code Python",   ORANGE),
    ("8",      "tables PostgreSQL",       BLEU_CLAIR),
]

for i, (val, label, color) in enumerate(kpis):
    row, col = divmod(i, 4)
    x = 0.3 + col * 3.25
    y = 1.55 + row * 1.8
    add_rounded_rect(s, Inches(x), Inches(y), Inches(3.1), Inches(1.65),
                     fill=RGBColor(0x2D, 0x4B, 0xA4))
    add_text(s, Inches(x), Inches(y + 0.1), Inches(3.1), Inches(1.0),
             val, size=40, bold=True, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(x), Inches(y + 1.1), Inches(3.1), Inches(0.5),
             label, size=12, color=BLANC,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.5), Inches(7.18), Inches(12.3), Inches(0.3),
         "JumiaPrix CI — ENSEA AS Data Science                     15 / 16",
         size=10, color=BLEU_CLAIR, align=PP_ALIGN.CENTER)

add_notes(s, "Ces chiffres parlent d'eux-mêmes. "
             "Près de 2000 produits, plus de 13000 snapshots de prix historisés, "
             "1088 promotions actives détectées automatiquement. "
             "Côté performance, 2500+ requêtes API servies avec ZÉRO erreur serveur. "
             "Côté qualité : 42 tests pytest, 12 services Docker orchestrés, "
             "plus de 5000 lignes de Python. "
             "C'est un projet de taille réelle, avec des données réelles, "
             "qui tourne en production sur ma machine.")


# ────────────────────────────────────────────────────────────
# SLIDE 16 — Conclusion
# ────────────────────────────────────────────────────────────
s = add_blank_slide()
set_bg(s, BLANC)
add_title_bar(s, "Conclusion", "Ce que j'ai construit et appris")

# Acquis
add_rounded_rect(s, Inches(0.5), Inches(1.6), Inches(6.1), Inches(2.6),
                 fill=BLANC, line=VERT)
add_rect(s, Inches(0.5), Inches(1.6), Inches(6.1), Inches(0.5), fill=VERT)
add_text(s, Inches(0.5), Inches(1.6), Inches(6.1), Inches(0.5),
         "Compétences acquises", size=14, bold=True, color=BLANC,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
acquis = [
    "Maîtrise de Scrapy et contournement Cloudflare",
    "Orchestration Airflow + Celery en production",
    "Conception DB relationnelle avancée (vues, triggers)",
    "API REST complète avec auth, tests, Swagger",
    "Intégration LLM avec prompt engineering",
    "DevOps : Docker, monitoring, CI/CD",
]
add_bullets(s, Inches(0.7), Inches(2.25), Inches(5.8), Inches(1.9),
            acquis, size=12)

# Perspectives
add_rounded_rect(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(2.6),
                 fill=BLANC, line=ORANGE)
add_rect(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.5), fill=ORANGE)
add_text(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.5),
         "Perspectives d'évolution", size=14, bold=True, color=BLANC,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
perspectives = [
    "Ajouter d'autres sources (Afrimalin, Abidjan.net)",
    "Prédiction des baisses de prix (modèle ML)",
    "Application mobile Flutter",
    "Déploiement cloud (AWS ou Hetzner)",
    "Notifications push via WhatsApp API",
    "Extension Chrome pour comparer depuis Jumia",
]
add_bullets(s, Inches(7.0), Inches(2.25), Inches(5.7), Inches(1.9),
            perspectives, size=12)

# Merci
add_rect(s, Inches(0), Inches(4.8), Inches(13.333), Inches(2.05), fill=BLEU_NUIT)
add_rect(s, Inches(0), Inches(4.8), Inches(13.333), Inches(0.05), fill=ORANGE)
add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.8),
         "Merci pour votre attention",
         size=40, bold=True, color=BLANC,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.4),
         "github.com/nathankindo008-dot/webscraping-pipeline-comparateur-prix",
         size=13, color=ORANGE_PALE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.4),
         "Des questions ?",
         size=18, bold=True, color=ORANGE,
         align=PP_ALIGN.CENTER)

add_footer(s, 16)
add_notes(s, "En conclusion, ce projet m'a permis de monter en compétence sur beaucoup "
             "d'outils en même temps : Scrapy, Airflow, Celery, PostgreSQL avancé, "
             "intégration LLM, Docker, monitoring. "
             "Les perspectives sont nombreuses : ajouter des sources, prédire les baisses "
             "de prix avec un modèle ML, faire une app mobile, déployer dans le cloud. "
             "Merci pour votre attention, je suis disponible pour répondre à vos questions.")


# ────────────────────────────────────────────────────────────
# Sauvegarde
# ────────────────────────────────────────────────────────────
output_file = "presentation_jumiaprix.pptx"
prs.save(output_file)
print(f"[OK] {len(prs.slides)} slides générées dans {output_file}")
